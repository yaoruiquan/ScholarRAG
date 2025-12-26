"""
统一文档加载器
支持多种文件格式：PDF, DOCX, TXT, MD, PPTX
"""

import os
from pathlib import Path
from typing import List, Optional
from langchain_core.documents import Document


def load_pdf(file_path: str) -> List[Document]:
    """加载 PDF 文件"""
    import fitz  # PyMuPDF
    
    documents = []
    try:
        doc = fitz.open(file_path)
        for page_num, page in enumerate(doc):  # type: ignore
            text = page.get_text()
            if text.strip():
                documents.append(Document(
                    page_content=text,
                    metadata={
                        "source": os.path.basename(file_path),
                        "page": page_num + 1,
                        "file_type": "pdf"
                    }
                ))
        doc.close()
    except Exception as e:
        print(f"  ⚠️ PDF 加载失败 {file_path}: {e}")
    
    return documents


def load_docx(file_path: str) -> List[Document]:
    """加载 Word 文档"""
    from docx import Document as DocxDocument
    
    documents = []
    try:
        doc = DocxDocument(file_path)
        
        # 提取所有段落文本
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text)
        
        # 提取表格内容
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                if row_text:
                    full_text.append(f"[表格] {row_text}")
        
        if full_text:
            documents.append(Document(
                page_content="\n".join(full_text),
                metadata={
                    "source": os.path.basename(file_path),
                    "page": 1,
                    "file_type": "docx"
                }
            ))
    except Exception as e:
        print(f"  ⚠️ DOCX 加载失败 {file_path}: {e}")
    
    return documents


def load_txt(file_path: str) -> List[Document]:
    """加载纯文本文件"""
    documents = []
    try:
        # 尝试多种编码
        for encoding in ['utf-8', 'gbk', 'gb2312', 'latin-1']:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                break
            except UnicodeDecodeError:
                continue
        else:
            print(f"  ⚠️ 无法解码文件 {file_path}")
            return []
        
        if content.strip():
            documents.append(Document(
                page_content=content,
                metadata={
                    "source": os.path.basename(file_path),
                    "page": 1,
                    "file_type": "txt"
                }
            ))
    except Exception as e:
        print(f"  ⚠️ TXT 加载失败 {file_path}: {e}")
    
    return documents


def load_markdown(file_path: str) -> List[Document]:
    """加载 Markdown 文件"""
    return load_txt(file_path)  # 与 TXT 处理相同


def load_pptx(file_path: str) -> List[Document]:
    """加载 PowerPoint 文件"""
    from pptx import Presentation
    
    documents = []
    try:
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = getattr(shape, "text", "")
                    if text.strip():
                        slide_text.append(text)
            
            if slide_text:
                documents.append(Document(
                    page_content="\n".join(slide_text),
                    metadata={
                        "source": os.path.basename(file_path),
                        "page": slide_num + 1,
                        "file_type": "pptx"
                    }
                ))
    except Exception as e:
        print(f"  ⚠️ PPTX 加载失败 {file_path}: {e}")
    
    return documents


# 支持的文件格式映射
FILE_LOADERS = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".doc": load_docx,  # 尝试用 docx 加载
    ".txt": load_txt,
    ".md": load_markdown,
    ".markdown": load_markdown,
    ".pptx": load_pptx,
    ".ppt": load_pptx,  # 尝试用 pptx 加载
}

# 支持的文件扩展名列表
SUPPORTED_EXTENSIONS = list(FILE_LOADERS.keys())


def load_document(file_path: str) -> List[Document]:
    """
    加载单个文档
    
    Args:
        file_path: 文件路径
        
    Returns:
        Document 列表
    """
    ext = Path(file_path).suffix.lower()
    
    if ext not in FILE_LOADERS:
        print(f"  ⚠️ 不支持的文件格式: {ext}")
        return []
    
    loader = FILE_LOADERS[ext]
    return loader(file_path)


def load_documents_from_directory(
    directory: str, 
    extensions: Optional[List[str]] = None
) -> List[Document]:
    """
    从目录加载所有支持的文档
    
    Args:
        directory: 目录路径
        extensions: 指定要加载的扩展名列表，默认加载所有支持的格式
        
    Returns:
        Document 列表
    """
    if extensions is None:
        extensions = SUPPORTED_EXTENSIONS
    
    documents = []
    dir_path = Path(directory)
    
    if not dir_path.exists():
        print(f"⚠️ 目录不存在: {directory}")
        return []
    
    # 遍历目录中的所有文件
    for file_path in dir_path.iterdir():
        if file_path.is_file() and file_path.suffix.lower() in extensions:
            print(f"  📄 加载: {file_path.name}")
            docs = load_document(str(file_path))
            documents.extend(docs)
    
    print(f"✅ 共加载 {len(documents)} 个文档片段")
    return documents


def get_supported_extensions_str() -> str:
    """获取支持的文件扩展名字符串（用于 UI 显示）"""
    return ", ".join(SUPPORTED_EXTENSIONS)
